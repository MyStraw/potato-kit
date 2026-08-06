"""
발표자료 생성 — 랩 세미나 10분 (potato-kit 시운전)
실행: conda run --no-capture-output -n potato python make_slides.py
산출: reports/slides-titanic-20260806.pptx + reports/figs/slide_*.png

/slides 스킬 준수 사항
  - 발표 1분당 1장 상한 → 10분 = 9장
  - 제목이 곧 결론
  - 슬라이드용 그림은 글자를 키워 새로 그림 (보고서용 재사용 금지)
  - 본문 24pt 이상, 표 5행 이내
  - 발표자 노트 포함
  - 마지막 장은 "그래서 무엇을"
"""
import platform
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).parent
FIGS = ROOT / "reports/figs"
FIGS.mkdir(parents=True, exist_ok=True)

# 슬라이드용: 글자 크게, 고해상도
plt.rcParams.update({
    "font.family": {"Windows": "Malgun Gothic", "Darwin": "AppleGothic"}.get(
        platform.system(), "NanumGothic"),
    "axes.unicode_minus": False,
    "font.size": 15, "axes.titlesize": 17, "axes.labelsize": 15,
    "xtick.labelsize": 14, "ytick.labelsize": 14, "legend.fontsize": 14,
    "figure.dpi": 150,
})

df = pd.read_csv(ROOT / "reports/cv_v2_summary.csv")
real = df[~df["모델"].str.startswith("Baseline")].reset_index(drop=True)
SHORT = {"LogisticRegression": "LogReg", "RandomForest": "RF",
         "HistGradientBoosting": "HistGB"}
labels = [SHORT[m] for m in real["모델"]]

# ── 그림 1: 모델 간 차이가 시드 변동 안에 있다 ──────────────────
fig, ax = plt.subplots(figsize=(9, 4.6))
x = np.arange(len(real))
ax.errorbar(x, real["일반분할"], yerr=real["일반_시드변동"], fmt="o", ms=13,
            capsize=9, capthick=2.5, lw=2.5, color="#2c5aa0", label="StratifiedKFold")
lo, hi = real["일반분할"].min(), real["일반분할"].max()
ax.axhspan(lo, hi, color="#e0a030", alpha=0.22,
           label=f"모델 간 격차 {hi-lo:.4f}")
ax.set_xticks(x); ax.set_xticklabels(labels)
ax.set_ylabel("정확도"); ax.set_ylim(0.818, 0.840)
ax.set_title("모델 간 격차(0.0016) < 시드 간 변동(최대 0.0040)")
ax.legend(loc="lower right", framealpha=0.95)
ax.grid(axis="y", alpha=0.3)
fig.tight_layout(); fig.savefig(FIGS / "slide_1_no_diff.png", transparent=True); plt.close(fig)

# ── 그림 2: 누수는 HGB 에서만 ──────────────────────────────────
fig, ax = plt.subplots(figsize=(9, 4.6))
w = 0.36
ax.bar(x - w/2, real["일반분할"], w, yerr=real["일반_시드변동"], capsize=6,
       color="#5a8fd6", label="StratifiedKFold")
ax.bar(x + w/2, real["그룹분할"], w, yerr=real["그룹_시드변동"], capsize=6,
       color="#c9553d", label="StratifiedGroupKFold (Ticket)")
for i, r in real.iterrows():
    gap = r["누수기여"]
    sig = gap > 2 * max(r["일반_시드변동"], r["그룹_시드변동"])
    ax.annotate(f"{gap:+.4f}" + ("  *" if sig else ""),
                (i, max(r["일반분할"], r["그룹분할"]) + 0.006),
                ha="center", fontsize=15 if sig else 13,
                fontweight="bold" if sig else "normal",
                color="#c0392b" if sig else "#555")
ax.set_xticks(x); ax.set_xticklabels(labels)
ax.set_ylabel("정확도"); ax.set_ylim(0.78, 0.85)
ax.set_title("Ticket 누수는 부스팅에서만 실질적 (-0.0219 = 시드변동의 5.5배)")
ax.legend(loc="lower left", framealpha=0.95); ax.grid(axis="y", alpha=0.3)
fig.tight_layout(); fig.savefig(FIGS / "slide_2_leakage.png", transparent=True); plt.close(fig)

# ── 그림 3: 층화 건전성 진단 ───────────────────────────────────
fig, ax = plt.subplots(figsize=(9, 4.2))
names = ["StratifiedKFold\n(기준)", "GroupKFold\n(v1)", "StratifiedGroupKFold\n(v2)"]
vals = [0.0023, 0.0348, 0.0017]
cols = ["#7f8c8d", "#c0392b", "#27ae60"]
b = ax.bar(names, vals, color=cols, width=0.55)
for rect, v, base in zip(b, vals, [1, 15.1, 0.8]):
    ax.text(rect.get_x() + rect.get_width()/2, v + 0.0012,
            f"{v:.4f}\n({base:g}배)", ha="center", fontsize=15, fontweight="bold")
ax.set_ylabel("Dummy 분류기 폴드 간 std"); ax.set_ylim(0, 0.046)
ax.set_title("층화 건전성 진단 — 모델과 무관하게 분할 결함만 드러낸다")
ax.grid(axis="y", alpha=0.3)
fig.tight_layout(); fig.savefig(FIGS / "slide_3_strat.png", transparent=True); plt.close(fig)

# ── PPTX ───────────────────────────────────────────────────────
FONT = "맑은 고딕"      # 후배들은 Windows
NAVY, GRAY, RED = RGBColor(0x1F, 0x3A, 0x5F), RGBColor(0x55, 0x55, 0x55), RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)   # 16:9
BLANK = prs.slide_layouts[6]


def add(title, body=None, img=None, notes="", title_size=30, body_size=24):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.15))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.name, p.font.color.rgb = Pt(title_size), True, FONT, NAVY
    tb.text_frame.word_wrap = True

    top = Inches(1.65)
    if img:
        pic = s.shapes.add_picture(str(img), Inches(0.9), top, width=Inches(11.5))
        top = Emu(pic.top + pic.height + Inches(0.2))
    if body:
        h = prs.slide_height - top - Inches(0.4)
        if h > Inches(0.5):
            box = s.shapes.add_textbox(Inches(0.9), top, Inches(11.5), h)
            tf = box.text_frame; tf.word_wrap = True
            for i, line in enumerate(body):
                par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                lv = line.startswith("  ")
                par.text = ("• " + line.strip()) if not lv else ("– " + line.strip())
                par.level = 1 if lv else 0
                par.font.size = Pt(max(24, body_size) if not lv else max(20, body_size - 3))
                par.font.name = FONT
                par.font.color.rgb = GRAY if lv else RGBColor(0x22, 0x22, 0x22)
                par.space_after = Pt(9)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


# 1
s = add("potato-kit 워크플로 시운전", title_size=40)
sub = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.6))
tf = sub.text_frame; tf.word_wrap = True
for i, (t, sz, c, b) in enumerate([
    ("Titanic 데이터로 본 '검증 설계의 함정'", 30, NAVY, True),
    ("EDA → 실험 → 방법론 감사 → 보고서 전 과정 검증", 22, GRAY, False),
    ("2026-08-06 · 891행 공개 데이터 · scikit-learn 1.9.0", 18, GRAY, False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = t; p.font.size = Pt(sz); p.font.name = FONT
    p.font.color.rgb = c; p.font.bold = b; p.space_after = Pt(14)
s.notes_slide.notes_text_frame.text = (
    "오늘 발표는 모델 성능 자랑이 아니라, 연구 워크플로를 점검하면서 "
    "검증 설계에서 어떤 함정에 빠졌는지에 대한 이야기입니다.")

# 2
add("모델 성능이 아니라 '그 숫자를 믿을 수 있는가'를 물었다",
    ["세 가지 질문",
     "  모델 간 성능 차이를 통계적으로 구분할 수 있는가?",
     "  같은 티켓(일행)을 공유하는 구조가 누수를 일으키는가?",
     "  검증 설계 자체가 결과를 얼마나 좌우하는가?",
     "세 번째 질문은 계획에 없었다 — 실험 도중 발견해서 추가됐다"],
    notes="세 번째 질문이 오늘의 핵심입니다. 원래 계획에 없었는데 "
          "감사 과정에서 우리 실험 자체의 결함이 드러나면서 추가됐습니다.")

# 3
add("891행 데이터에 숨어 있던 구조 — 승객은 독립이 아니다",
    ["Kaggle Titanic train 891행 × 12열, 타깃 61.6% vs 38.4% (약한 불균형)",
     "결측: Cabin 77%, Age 20% → HasCabin 플래그 + Age 결측 지시자",
     "구조적 위험: 같은 Ticket = 같은 일행. 생사를 함께하는 경향",
     "  train 내부 Ticket 중복 210건, train–test 걸침 115건",
     "  → 폴드를 나눌 때 일행이 흩어지면 정보가 샌다 (EDA 단계의 사전 경고)"],
    notes="EDA에서 미리 경고했던 부분입니다. 승객 한 명 한 명이 독립이라고 "
          "가정하고 무작위로 폴드를 나누면, 같은 가족이 학습셋과 검증셋에 "
          "나뉘어 들어갑니다.")

# 4
add("방법 — 검증 설계를 v1에서 v2로 바꿨다",
    ["v1: StratifiedKFold vs GroupKFold(Ticket), 단일 시드 42",
     "v2: StratifiedKFold vs StratifiedGroupKFold(Ticket), 5개 시드",
     "공통: 전처리는 전부 Pipeline 내부 → 폴드 안에서만 fit (누수 차단)",
     "파생변수 Title·FamilySize·IsAlone·HasCabin 은 타깃 미사용 → 폴드 밖 생성 무해",
     "하이퍼파라미터는 튜닝하지 않은 관례적 기본값 (의도적)",
     "  → 따라서 '모델 계열의 우열'이 아니라 '임의의 한 점끼리 비교'다"],
    notes="전처리를 파이프라인 안에 넣은 것은 v1부터 지킨 원칙입니다. "
          "튜닝을 안 한 것은 의도적이고, 그래서 결론에서 모델 우열을 주장하지 않습니다.")

# 5
add("판정 기준 — '무엇에 대한 변동인가'를 구분했다",
    ["폴드 간 std: 같은 데이터를 5번 나눠 본 흩어짐. 신뢰구간이 아니다",
     "  폴드끼리 훈련셋이 겹쳐 독립이 아니므로 유의성 판정에 쓸 수 없다",
     "시드 간 변동: 무작위성의 원천을 바꿔가며 반복. v2는 이것으로 판단",
     "누수 기여 = 일반분할 − 그룹분할. 시드 변동의 2배 초과일 때만 인정",
     "층화 건전성: Dummy 분류기의 폴드 간 std — 모델과 무관한 리트머스지"],
    notes="v1은 폴드 간 std를 유의성 판정 임계값으로 썼는데 이건 통계적 근거가 "
          "없습니다. v2에서 시드 간 변동으로 바꿨습니다.")

# 6
add("결과 ① 세 모델은 구분되지 않는다", img=FIGS / "slide_1_no_diff.png",
    notes="격차 0.0016, 시드 변동 최대 0.0040. 격차가 변동보다 작으므로 "
          "순위표를 만들면 안 됩니다. 단일 시드로 끝냈다면 HistGB를 "
          "'최고 모델'로 잘못 결론지었을 겁니다.")

# 7
add("결과 ② 누수는 부스팅에서만 실질적이다", img=FIGS / "slide_2_leakage.png",
    notes="LR과 RF는 시드 변동 안이라 누수가 있다고 말할 수 없습니다. "
          "HistGB만 5.5배. 다만 원인은 분리하지 못했습니다 — 모델 특성인지, "
          "부스팅에 선형모델용 전처리를 강제한 불이익인지, 과적합인지.")

# 8
add("결과 ③ 검증 설계가 결과를 만든다", img=FIGS / "slide_3_strat.png",
    notes="GroupKFold는 층화를 하지 않습니다. 그래서 v1이 잰 '누수'에는 "
          "층화 상실과 단일분할 노이즈가 섞여 실제보다 약 1.4배 부풀려졌습니다. "
          "−0.0303 → −0.0219. 진단 지표는 Dummy 분류기의 폴드 간 std입니다.")

# 9
add("그래서 — 검증 설계를 먼저 감사하고 시작한다",
    ["바로 쓸 수 있는 체크",
     "  개체 그룹이 있으면 StratifiedGroupKFold. GroupKFold는 층화가 없다",
     "  단일 시드 순위를 믿지 않는다. 최소 5시드 + 시드 간 변동 병기",
     "  Dummy 분류기 폴드 std 로 층화 건전성을 30초 만에 진단",
     "다음 단계",
     "  모델별 적합 전처리로 재실험 → 누수 원인 분해 (전처리 vs 모델 특성)",
     "  성(姓) 정보로 그룹 키 보강 → 현재 추정치는 누수의 하한"],
    notes="오늘 얻은 것은 타이타닉 성능이 아니라 재사용 가능한 검증 체크리스트입니다. "
          "이 세 줄은 어느 프로젝트에도 그대로 적용됩니다.")

out = ROOT / "reports/slides-titanic-20260806.pptx"
prs.save(out)

print(f"✓ {out.name} — {len(prs.slides.__iter__.__self__._sldIdLst)}장")
print(f"  슬라이드용 그림 3개 → reports/figs/slide_*.png")
notes_n = sum(1 for s in prs.slides if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())
print(f"  발표자 노트 {notes_n}장에 포함")
print(f"  크기 {out.stat().st_size/1024:.0f} KB")
